import ctypes
import time
from tkinter import filedialog
from tkinter import * 
from tkinter.ttk import *
import networkx as nx
import re
import os
import math

import tkinter.font as tf
from numpy.ma import var
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO_SHAPE
from pptx.util import Inches,Pt
from pptx.enum.shapes import MSO_CONNECTOR_TYPE as MSO_CONNECTOR
import tkinter.colorchooser as cc
from tkinter import messagebox
import random
import tkinter

'''
我把文件的读取整体都放到了选择文件路径的按钮命令里
'''


#全局变量
LAYOUT_KIND = 0  # 0-spring 1-circular 2-random 3-shell 4-spectral
CONNECTOR_CHOICE = 0 #0-直线； 1-曲线； 2-直角曲线
NODE_SIZE_SELECTION = 0
NODE_COLOR_SELECTION = 1
file_path = ''
attr_file_path = ''
NODE_MAX_SIZE = 0.6  #Inches
NODE_MIN_SIZE = 0.3 #Inches
LINE_SIZE = 0.1  #Inches
FONT_SIZE = 15   #Pt
indexx = 0
edges = []# 保存边关系，（序号1，序号2）
colors = []#颜色列表

#节点、连线、字体的颜色
#NODE_COLOR_RGB = RGBColor(64,128,128)
LINE_COLOR_RGB = RGBColor(192,192,192)
FONT_COLOR_RGB = RGBColor(255,255,255)
n = dict() # 保存节点,(序号，[节点名称,属性，度])
n_reverse = dict() # (节点名称，序号)
attribute = dict()
index_of_attribute = list()#序号，属性名称

#选择布局类型
def getposition(G,choice):
    if choice ==1:
        return nx.circular_layout(G)
    elif choice ==2:
        return nx.random_layout(G)
    elif choice ==3:
        return nx.shell_layout(G)
    elif choice ==4:
        return nx.spectral_layout(G)
    else :
        return nx.spring_layout(G)

def get_node_size(n_degree,dif_degree,dif_size):
    temp = n_degree/100
    node_size = NODE_MIN_SIZE + (dif_size) * temp
    if dif_degree!=0:
        node_size += temp*10/ (dif_degree)
    return node_size

#根据节点数量随机生成颜色, input:int output:list
def get_colors(n):
    colors = list(map(lambda i: "#" + "%06x" % random.randint(0, 0xFFFFFF),range(n)))
    return colors

mark = int()
def change_node_color(str):
    global mark
    global colors
    choose = cc.askcolor()
    if choose[1] is None:
        return
    else:
        i = 0
        count = 0
        for key in attribute:
            if str == key:
                i = count
            else:
                count += 1
        tmp = choose[1]
        colors[i] = tmp
        color = colors[i]
        labels[i].configure(background=color,text='属性'+str+'颜色'+color)

labels = []
def color_selection():
    top = Toplevel(window)
    screenwidth = window.winfo_screenwidth()
    screenheight = window.winfo_screenheight()
    top.geometry('%dx%d+%d+%d' % (width, height, (screenwidth - width) / 2, (screenheight - height) / 2))
    #top.attributes("-toolwindow", 1)
    #top.wm_attributes("-topmost", 1)
    #buttons = []
    global labels
    global colors
    global index_of_attribute
    global attribute
    #n 保存节点,(序号，[节点名称,属性，度])
    label_combo = Label(master=top, text = '请选择需要修改颜色的属性')
    label_combo.grid(row=1, column=1)
    combo_top = Combobox(top)
    tmp = []
    for key in attribute:
        tmp.append(key)
    tuple_att = tuple(tmp)
    combo_top['values'] = tuple_att
    combo_top.grid(row=1, column=2)
    set(combo_top['value'][0])
    for i in range(len(colors)):
        label = Label(master=top)
        color = colors[i]
        label.configure(background=color, text="属性"+index_of_attribute[i]+"默认颜色")
        label.grid(row=(i)%25, column=int((i)/25))
        labels.append(label)
    '''for i in range(len(attribute)):
        #button = Button(text="选择"+index_of_attribute[i]+"属性颜色", master=window,command = change_node_color(colors,i))
        button = Button(text=index_of_attribute[i], master=top,command = lambda :change_node_color(button['text']))
        button.grid(row=i+1, column=2)
        buttons.append(button)'''
    '''for i in range(len(attribute)):
        color = colors[i]
        labels[i].configure(background=color,text='颜色'+color)'''
    e1 = combo_top
    button = Button(text="选择属性颜色", master=top,command = lambda :change_node_color(e1.get()))
    button.grid(row=3, column=5)
    button_off = Button(text="确定", master=top,command=top.destroy)
    button_off.grid(row=3, column=6)
    window.wait_window(top)
    return colors

def change_node_color_without_attr(str):
    global mark
    global colors
    choose = cc.askcolor()
    if choose[1] is None:
        return
    else:
        i = 0
        count = 0
        global n_reverse
        i = n_reverse[str]
        tmp = choose[1]
        colors[i] = tmp
        color = colors[i]
        labels[i].configure(background=color,text=n[i][0]+'节点颜色'+color)

def color_selection_without_attribute():
    top = Toplevel(window)
    screenwidth = window.winfo_screenwidth()
    screenheight = window.winfo_screenheight()
    top.geometry('%dx%d+%d+%d' % (width, height, (screenwidth - width) / 2, (screenheight - height) / 2))
    global labels
    global colors
    global n
    #n 保存节点,(序号，[节点名称,属性，度])
    label_combo = Label(master=top, text = '请选择需要修改颜色的节点')
    label_combo.grid(row=1, column=5)
    combo_top = Combobox(top)
    tmp = []
    for i in range(len(n)):
        tmp.append(n[i][0])
    tuple_att = tuple(tmp)
    combo_top['values'] = tuple_att
    #combo_top.current(n[0][0])
    combo_top.grid(row=1, column=6)
    set(combo_top['value'][0])

    for i in range(1,len(n)+1):
        label = Label(master=top)
        color = colors[i-1]
        label.configure(background=color,text="节点"+n[i-1][0]+"默认颜色")
        label.grid(row=(i-1)%25, column=int((i-1)/25))
        labels.append(label)
    '''for i in range(len(attribute)):
        #button = Button(text="选择"+index_of_attribute[i]+"属性颜色", master=window,command = change_node_color(colors,i))
        button = Button(text=index_of_attribute[i], master=top,command = lambda :change_node_color(button['text']))
        button.grid(row=i+1, column=2)
        buttons.append(button)'''
    '''for i in range(len(attribute)):
        color = colors[i]
        labels[i].configure(background=color,text='颜色'+color)'''
    e1 = combo_top
    button = Button(text="选择节点颜色", master=top,command=lambda:change_node_color_without_attr(e1.get()))
    button.grid(row=1, column=7)
    button_off = Button(text="确定", master=top,command=top.destroy)
    button_off.grid(row=1, column=8)
    window.wait_window(top)
    return colors

# 选择连接的节点的哪个端点（0-上，1-左，2-上，3-右）
def getdirec(node1, node2):
    direc = []
    temp_x = node2[0] - node1[0]
    temp_y = node2[1] - node1[1]
    if (temp_y >= 0 and abs(temp_x) < temp_y):
        direc.append(0)
        direc.append(2)
    elif (temp_y <= 0 and abs(temp_x) < abs(temp_y)):
        direc.append(2)
        direc.append(0)
    elif (temp_x >= 0 and abs(temp_y) < temp_x):
        direc.append(3)
        direc.append(1)
    else:
        direc.append(1)
        direc.append(3)
    return direc

#主要功能
def main_function():

    global LAYOUT_KIND
    #global CONNECTOR_KIND
    global NODE_SIZE_SELECTION

    #CONNECTOR_KIND = combo1.get()
    #LAYOUT_KIND = combo.get()
    #NODE_SIZE_SELECTION = combo2.get()


    if attr_file_path == '':
        G = nx.Graph()
        G.add_edges_from(edges)
        G.add_nodes_from(list(range(len(n))))
        #LAYOUT_KIND = combo.get()
        #按照选择的布局类型获取点的布局
        pos=getposition(G,LAYOUT_KIND)

        #生成幻灯片
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes

        nodes=[]
        #unit_size = (NODE_MAX_SIZE-NODE_MIN_SIZE)/(max_degree-min_degree)
        global colors
        global NODE_COLOR_SELECTION
        #NODE_COLOR_SELECTION = combo_color.get()
        #生成节点个数个随机颜色
        colors = get_colors(len(n))
        if NODE_COLOR_SELECTION.get() == 0:
            colors = color_selection_without_attribute()
        #NODE_SIZE_SELECTION = combo2.get()
        #生成节点，插入到shapes中
        for i in range(len(n)):
            #default:  width-10 inches   height-7.46inches
            (x, y) = pos[i]
            #node_size = NODE_MIN_SIZE + unit_size * (n[i][2]-min_degree)
            #node_size = NODE_MIN_SIZE + (NODE_MAX_SIZE - NODE_MIN_SIZE) * math.sin(n[i][2] / 100)
            #node_size = NODE_MIN_SIZE + (NODE_MAX_SIZE - NODE_MIN_SIZE) * (n[i][2] / 100) + (n[i][2]/100)/ (max_degree - min_degree)
            #print(NODE_SIZE_SELECTION.get())
            if NODE_SIZE_SELECTION.get() == 0:
                #print('.......')
                node_size = get_node_size(n[i][2], dif_degree, dif_size)
            elif NODE_SIZE_SELECTION.get()==1:
                node_size = 0.5
                #print('******')
            nodes.append(shapes.add_shape(MSO_SHAPE.OVAL, Inches(5*(x+1)), Inches(3.7*(1-y)), Inches(node_size),Inches(node_size)))
            frame = nodes[i].text_frame
            frame.word_wrap = FALSE
            para = frame.paragraphs[0]
            run = para.add_run()
            run.text = n[i][0]
            run.font.size = Pt(FONT_SIZE)
            #n[i][1] 存的是属性，目前是字符类型
            run.font.color.rgb = FONT_COLOR_RGB
            #set the color of the cycle
            fill = nodes[i].fill
            fill.solid()
            att = n[i][1]
            color = Hex_to_RGB(colors[i])

            fill.fore_color.rgb = RGBColor(color[0],color[1],color[2])
            nodes[i].line.color.rgb = RGBColor(color[0],color[1],color[2])
    else:
         #新建图对象
        G = nx.Graph()

        G.add_edges_from(edges)
        G.add_nodes_from(list(range(len(n))))
        #LAYOUT_KIND = combo.get()
        #按照选择的布局类型获取点的布局
        pos=getposition(G,LAYOUT_KIND)

        #生成幻灯片
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes

        nodes=[]
        #unit_size = (NODE_MAX_SIZE-NODE_MIN_SIZE)/(max_degree-min_degree)

        #生成种类个数个随机颜色
        colors = get_colors(len(attribute))
        #atrribute:属性字典 key:属性 Value:颜色RGB
        i = 0


        #NODE_COLOR_SELECTION = combo_color.get()
        if NODE_COLOR_SELECTION.get() == 0:
            colors = color_selection()
        for key in attribute:
            color = Hex_to_RGB(colors[i])
            #color_RGB = RGBColor(color[0],color[1],color[2])
            attribute[key] = color
            i += 1

        #NODE_SIZE_SELECTION = combo2.get()
        #生成节点，插入到shapes中
        for i in range(len(n)):
            #default:  width-10 inches   height-7.46inches
            (x, y) = pos[i]
            if NODE_SIZE_SELECTION.get() == 0:
                node_size = get_node_size(n[i][2],dif_degree,dif_size)
                #print('......')
            elif NODE_SIZE_SELECTION.get() == 1:
                node_size=0.4
                #print('******')
            nodes.append(shapes.add_shape(MSO_SHAPE.OVAL, Inches(5*(x+1)), Inches(3.7*(1-y)), Inches(node_size),Inches(node_size)))
            frame = nodes[i].text_frame
            frame.word_wrap = FALSE
            para = frame.paragraphs[0]
            run = para.add_run()
            run.text = n[i][0]
            run.font.size = Pt(FONT_SIZE)
            #n[i][1] 存的是属性，目前是字符类型
            run.font.color.rgb = FONT_COLOR_RGB
            #set the color of the cycle
            fill = nodes[i].fill
            fill.solid()
            att = n[i][1]
            color = attribute[att]
            fill.fore_color.rgb = RGBColor(color[0],color[1],color[2])
            nodes[i].line.color.rgb = RGBColor(color[0],color[1],color[2])



    #生成节点之间的连接线，插入到shapes中
    for i in range(len(G.edges)):

        #CONNECTOR_CHOICE = combo1.get()
        if int(CONNECTOR_CHOICE.get()) == 1:
            connector = shapes.add_connector(
            MSO_CONNECTOR.CURVE, 0, 0, 0, 0
        )
        elif int(CONNECTOR_CHOICE.get()) == 2:
            connector = shapes.add_connector(
            MSO_CONNECTOR.ELBOW, 0, 0, 0, 0
        )
        else:
            connector = shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, 0, 0, 0, 0
        )

        direc = getdirec(pos[edges[i][0]],pos[edges[i][1]])
        connector.begin_connect(nodes[edges[i][0]],direc[0])
        connector.end_connect(nodes[edges[i][1]],direc[1])
        connector.line.color.rgb = LINE_COLOR_RGB

    #把节点位置放到连接线形状的后面，这样节点可以先是在连接线上方
    for i in range(len(nodes)):
        cursor_sp = shapes[len(shapes)-1]._element
        cursor_sp.addnext(shapes[0]._element)

    #保存文件，将文件输入路径进行分割，取.txt前面的部分
    file_save_path = re.split("\.txt", file_path)
    save_path = file_save_path[0]
    localtime = time.localtime(time.time())
    flag = str(localtime.tm_min) + '-'+str(localtime.tm_sec)
    prs.save(save_path+flag+'.pptx')
    messagebox.showinfo("保存", "文件保存完毕！")
    os.startfile(save_path+flag+'.pptx')
    window.destroy()



#分割线：后面都是程序界面的代码
#浏览文件
def browse_button():
    filename = filedialog.askopenfilename()
    if 'txt' in filename:
        global file_path
        file_path = filename
        file = open(file_path, 'rb')
        line = file.readline()
        line_de = line.decode()
        if 'net' not in line_de:
            messagebox.showinfo("错误", "网络关系文件输入错误！")
        else:
            lbl1.configure(text = '关系信息文件：'+filename)
            file = open(file_path,'rb')
            line = file.readline()
            global eges

            # 方便快速由序号找到节点名，由节点名找到对应的序号
            global n # 保存节点,(序号，[节点名称,属性，度])
            global n_reverse# (节点名称，序号)

            # 读取文件中的节点和边
            index = 0 # 记录序号

            line = file.readline()
            while line:
                #edge = line.decode().split('，')
                line_de = line.decode()
                edge = re.split(r'[，,\r]', line_de)
                x = edge[0]
                x_index,y_index=0,0
                if x in n_reverse:
                    x_index = n_reverse.get(x)
                else:
                    n[index]= [x, None, 0]
                    n_reverse[x]=index
                    x_index = index
                    index += 1
                # 节点的度+1
                temp = n[x_index]
                temp[2] += 1
                y = edge[1]
                if y in n_reverse:
                    y_index = n_reverse.get(y)
                else:
                    n[index]=[y, None, 0]
                    n_reverse[y]=index
                    y_index = index
                    index += 1
                # 节点的度+1
                temp = n[y_index]
                temp[2] += 1
                edges.append((x_index,y_index))
                line = file.readline()
            global indexx
            indexx = index
            max_degree = 0
            min_degree = 100
            for i in range(len(n)):
                if (n[i][2] > max_degree):
                    max_degree = n[i][2]
                if (n[i][2] < min_degree):
                    min_degree = n[i][2]
            global dif_degree
            global dif_size
            dif_degree = max_degree - min_degree
            dif_size = NODE_MAX_SIZE - NODE_MIN_SIZE
            if len(n) > 100:
                messagebox.showinfo("错误", "节点数量应该小于100！")
                return
    else:
        messagebox.showinfo("错误", "错误的文件格式输入")


#上传属性文件
def browse_attr_button():
    attr_filename = filedialog.askopenfilename()
    if 'txt' in attr_filename:
        global attr_file_path
        attr_file_path = attr_filename
        attr_file = open(attr_file_path, 'rb')
        attr_line = attr_file.readline()
        line_de = attr_line.decode()
        if 'node' not in line_de:
            messagebox.showinfo("错误", "节点文件输入错误！")
        else:
            attr_lbl.configure(text = '属性信息文件：' +attr_filename)
            global attribute #保存属性
            index_attribute = 0
            attr_file = open(attr_file_path,'rb')
            attr_line = attr_file.readline()
            line_de = attr_line.decode()
            attr_line = attr_file.readline()
            global index_of_attribute
            global indexx
            while attr_line:
                line_de = attr_line.decode()
                attr = re.split(r'[，,\r]', line_de)
                #读取属性，保存在属性字典中
                if attr[1] not in attribute:
                    attribute[attr[1]] = ''
                    index_of_attribute.append(attr[1])
                    index_attribute += 1
                if attr[0] not in n_reverse:
                    n[indexx]=[attr[0], attr[1], 0]
                    n_reverse[attr[0]]=indexx
                    indexx += 1
                else:
                    i = n_reverse.get(attr[0])
                    temp = n[i]
                    temp[1] = attr[1]
                attr_line = attr_file.readline()
    else:
        messagebox.showinfo("错误", "错误的文件格式输入")

#生成ppt的按钮
def generate_button():
    main_function()

#节点颜色选择命令执行
def choose_node_color():
    choose = cc.askcolor()
    global NODE_COLOR_RGB
    tmp = choose[1]
    tmp = Hex_to_RGB(tmp)
    NODE_COLOR_RGB = RGBColor(tmp[0], tmp[1], tmp[2])
    #btn_node_color.config(text=tmp)
    #lbl_node_color.configure(background=choose[1],text=tmp)

#线条颜色选择命令执行
def choose_line_color():
    choose = cc.askcolor()
    if choose[1] is None:
        return
    global LINE_COLOR_RGB
    tmp = choose[1]
    tmp = Hex_to_RGB(tmp)
    LINE_COLOR_RGB = RGBColor(tmp[0], tmp[1], tmp[2])
    #btn_line_color.config(text=tmp)
    lbl_line_color.configure(background=choose[1],text=tmp)


#字体颜色选择命令执行
def choose_font_color():
    choose = cc.askcolor()
    if choose[1] is None:
        return
    global FONT_COLOR_RGB
    tmp = choose[1]
    tmp = Hex_to_RGB(tmp)
    FONT_COLOR_RGB = RGBColor(tmp[0], tmp[1], tmp[2])
    #btn_font_color.config(text=tmp)
    lbl_font_color.configure(background=choose[1],text=tmp)

# 16进制颜色格式颜色转换为RGB格式
def Hex_to_RGB(hex):
    r = int(hex[1:3],16)
    g = int(hex[3:5],16)
    b = int(hex[5:7], 16)
    rgb = []
    rgb.append(r)
    rgb.append(g)
    rgb.append(b)
    return rgb


window = Tk()
window.title('PPT Net Maker')
width = 600
height = 500
ft = tf.Font(family='Times', size=30, weight='bold',slant='italic')
ft1 = tf.Font(family='黑体', size=10)
screenwidth = window.winfo_screenwidth()
screenheight = window.winfo_screenheight()
window.geometry('%dx%d+%d+%d'%(width, height, (screenwidth-width)/2, (screenheight-height)/2))


folder_path = StringVar()

title = Label(master = window , text = 'PPT Net Maker',font = ft)
#title.grid(row=1,column=1)

title.place(x=150, y=20)

lbl1 = Label(master=window, text = '上传关系信息文件(.txt):',font = ft1,wraplength = 230)
lbl1.place(x = 80, y = 100)

button2 = Button(text="Browse", command=browse_button)
button2.place(x=350, y=95)

#选择属性文件
attr_lbl = Label(master=window, text = '上传属性信息文件(.txt):\n(若不需要此项忽略)',font = ft1,wraplength = 230)
attr_lbl.place(x=80, y=130)

attr_button = Button(text="Browse", command=browse_attr_button)
attr_button.place(x=350, y=135)

#node_size_selection = Label(master=window, text = '是否自动调整节点大小：\n0-是, 1-否\n')
#node_size_selection.grid(row=5, column=1)
#combo2 = Combobox(window)
#combo2['values'] = (0,1)
#combo2.current(0)
#combo2.grid(row=5, column=2)
#NODE_SIZE_SELECTION = combo2.get()

NODE_SIZE_SELECTION = IntVar()
node_size_selection = Label(master=window, text = '是否自动调整节点大小:\n',font = ft1 )
node_size_selection.place(x = 80, y = 180 )
ns_default1 = Radiobutton(window, text='是', value=0, variable=NODE_SIZE_SELECTION)
ns_default1.place(x = 300,y=180)
ns_default2 = Radiobutton(window, text='否', value=1, variable=NODE_SIZE_SELECTION)
ns_default2.place(x = 350 ,y = 180)
NODE_SIZE_SELECTION.set(0)

#node_color_selection = Label(master=window, text = '是否手动调整节点颜色:\n0-是, 1-否\n')
#node_color_selection.grid(row=6, column=1)
#combo_color = Combobox(window)
#combo_color['values'] = (0,1)
#combo_color.current(1)
#combo_color.grid(row=6, column=2)
#NODE_COLOR_SELECTION = combo_color.get()
NODE_COLOR_SELECTION = IntVar()
node_color_selection = Label(master=window, text = '是否手动调整节点颜色:\n',font = ft1 )
node_color_selection.place(x = 80, y = 210 )
nc_default1 = Radiobutton(window, text='是', value=0, variable=NODE_COLOR_SELECTION)
nc_default1.place(x = 300,y=210)
nc_default2 = Radiobutton(window, text='否', value=1, variable=NODE_COLOR_SELECTION)
nc_default2.place(x = 350 ,y = 210)
NODE_COLOR_SELECTION.set(1)

#lbl2 = Label(master=window, text = '选择节点排布样式：\n0-spring_layout\n1-circular_layout\n2-random_layout\n3-shell_layout\n4-spectral_layout')
#lbl2.grid(row=7, column=1)
#combo = Combobox(window)
#combo['values'] = (0,1,2,3,4)
#combo.current(0)
#combo.grid(row=7, column=2)
#LAYOUT_KIND = combo.get()
LAYOUT_KIND = IntVar()
lbl2 = Label(master=window, text = '选择节点布局样式:\n',font = ft1 )
lbl2.place(x = 80, y = 250 )
lbl2_default1 = Radiobutton(window, text='弹簧布局', value=0, variable=LAYOUT_KIND)
lbl2_default1.place(x = 300,y=240)
lbl2_default2 = Radiobutton(window, text='圆形布局', value=1, variable=LAYOUT_KIND)
lbl2_default2.place(x = 380 ,y = 240)
lbl2_default3 = Radiobutton(window, text='随机布局', value=2, variable=LAYOUT_KIND)
lbl2_default3.place(x = 460 ,y = 240)
lbl2_default4 = Radiobutton(window, text='壳牌布局', value=3, variable=LAYOUT_KIND)
lbl2_default4.place(x = 300 ,y = 260)
lbl2_default5 = Radiobutton(window, text='光谱布局', value=4, variable=LAYOUT_KIND)
lbl2_default5.place(x = 380 ,y = 260)
LAYOUT_KIND.set(0)

#lbl3 = Label(master=window, text = '选择连接线样式：\n0-直线\n1-曲线\n2-直角曲线')
#lbl3.grid(row=9, column=1)
#combo1 = Combobox(window)
#combo1['values'] = (0,1,2)
#combo1.current(0)
#combo1.grid(row=9, column=2)


CONNECTOR_CHOICE = IntVar()
lbl3 = Label(master=window, text = '选择连接线样式:\n',font = ft1 )
lbl3.place(x = 80, y = 290 )
lbl3_default1 = Radiobutton(window, text='直线', value=0, variable=CONNECTOR_CHOICE)
lbl3_default1.place(x = 300,y=290)
lbl3_default2 = Radiobutton(window, text='曲线', value=1, variable=CONNECTOR_CHOICE)
lbl3_default2.place(x = 350 ,y = 290)
lbl3_default3 = Radiobutton(window, text='直角曲线', value=2, variable=CONNECTOR_CHOICE)
lbl3_default3.place(x = 400 ,y = 290)
CONNECTOR_CHOICE.set(0)



'''节点按钮
btn_node_color = Button(text='选择节点颜色', command=choose_node_color)
btn_node_color.grid(row=6, column=2)
lbl_node_color = Label(master=window)
lbl_node_color.configure(background="#408080",text='默认颜色为墨绿色')
lbl_node_color.grid(row=6, column=1)'''

#线条按钮
btn_line_color = Button(text='选择线条颜色', command=choose_line_color)
btn_line_color.place( x=170 , y=330)
lbl_line_color = Label(master=window)
lbl_line_color.configure(background="#C0C0C0",text='默认颜色为灰色')
lbl_line_color.place(x=80, y=333)

#字体按钮
btn_font_color = Button(text='选择字体颜色', command=choose_font_color)
btn_font_color.place(x=390, y=330)
lbl_font_color = Label(master=window)
lbl_font_color.configure(background="white",text='默认颜色为白色')
lbl_font_color.place(x=300, y=333)



button3 = Button(text="保存文件",command=generate_button)
button3.place(x=250, y=400,width=80,height=30)

mainloop()
